"""Design tools — branded diagrams, SVG graphics, HTML mockups, designed PDFs, slides.

All visual output follows the deployment's brand, loaded from brand.yaml
(overlay wins over the built-in kbots placeholder brand in config/).
Rendering runs in headless Chromium (Playwright, already a core dependency);
slide decks use python-pptx (uv sync --extra design).

Typical uses:
- render_diagram: architecture/flow/sequence diagrams from Mermaid text
- render_svg: hand-authored vector graphics (badges, cards, logos)
- render_html: screenshot mockups (landing pages, emails, UI ideas)
- html_to_pdf: designed PDF reports — full CSS, brand header/footer
- create_slides: markdown outline → PowerPoint deck

Diagram/HTML rendering loads Pico.css / Tailwind / mermaid.js from jsDelivr,
so the host needs outbound network access.
"""

import base64
import ipaddress
import json
import logging
import os
import re
import socket
import xml.etree.ElementTree as ET
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from pathlib import Path

import yaml

from src.core.base import KBOTS_TMP, PROJECT_ROOT, ToolContext, resolve_config_file
from src.core.tools import tool
from src.tools.ingest import _BLOCKED_NETS, validate_file_path

logger = logging.getLogger(__name__)

_PW_HINT = "Error: playwright not installed. Run: uv add playwright && playwright install chromium"
_PPTX_HINT = "python-pptx not installed. Run: uv sync --extra design (then restart the service)"

PAGE_SIZES = ["A4", "A3", "A5", "Letter", "Legal", "Tabloid"]
HTML_STYLES = ["clean", "tailwind", "none"]

_DEFAULT_BRAND = {
    "name": "kbots",
    "tagline": "",
    "logo": "",
    "colors": {
        "primary": "#4F46E5", "secondary": "#334155", "accent": "#14B8A6",
        "background": "#FAFAFA", "surface": "#FFFFFF",
        "text": "#1E293B", "muted": "#64748B",
    },
    "font_family": "Inter, -apple-system, 'Segoe UI', Roboto, sans-serif",
    "font_url": "",
    "footer_text": "",
}


def _load_brand() -> dict:
    """Load brand config (overlay brand.yaml wins over core's placeholder)."""
    brand = {k: (dict(v) if isinstance(v, dict) else v) for k, v in _DEFAULT_BRAND.items()}
    try:
        path = resolve_config_file("brand.yaml")
        if path.is_file():
            loaded = yaml.safe_load(path.read_text()) or {}
            for key, value in loaded.items():
                if key == "colors" and isinstance(value, dict):
                    brand["colors"].update(value)
                elif key in brand:
                    brand[key] = value
    except Exception as e:
        logger.warning(f"brand.yaml could not be loaded, using defaults: {e}")
    return brand


def _brand_tokens_line(brand: dict) -> str:
    c = brand["colors"]
    return (f"Brand tokens: primary {c['primary']}, secondary {c['secondary']}, "
            f"accent {c['accent']}, text {c['text']}, muted {c['muted']} "
            f"(CSS: var(--brand-primary) etc.)")


def _brand_css_vars(brand: dict) -> str:
    c = brand["colors"]
    return (
        ":root{"
        f"--brand-primary:{c['primary']};--brand-secondary:{c['secondary']};"
        f"--brand-accent:{c['accent']};--brand-background:{c['background']};"
        f"--brand-surface:{c['surface']};--brand-text:{c['text']};"
        f"--brand-muted:{c['muted']};"
        "}"
    )


def _brand_head(brand: dict, style: str) -> str:
    """Build the <head> content that applies the brand + chosen style system."""
    head = ['<meta charset="utf-8">',
            '<meta name="viewport" content="width=device-width, initial-scale=1">']
    if brand.get("font_url"):
        head.append(f'<link rel="stylesheet" href="{brand["font_url"]}">')
    if style == "clean":
        head.append('<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/@picocss/pico@2/css/pico.min.css">')
        c = brand["colors"]
        head.append(
            "<style>"
            f":root{{--pico-primary:{c['primary']};--pico-primary-background:{c['primary']};"
            f"--pico-background-color:{c['background']};--pico-color:{c['text']};"
            f"--pico-muted-color:{c['muted']};}}"
            f"body{{font-family:{brand['font_family']};}}"
            "</style>"
        )
    elif style == "tailwind":
        c = brand["colors"]
        head.append('<script src="https://cdn.jsdelivr.net/npm/@tailwindcss/browser@4"></script>')
        head.append(
            '<style type="text/tailwindcss">'
            f"@theme {{ --color-brand: {c['primary']}; --color-brand-secondary: {c['secondary']}; "
            f"--color-brand-accent: {c['accent']}; --font-sans: {brand['font_family']}; }}"
            "</style>"
        )
    head.append(f"<style>{_brand_css_vars(brand)}</style>")
    return "\n".join(head)


def _resolve_asset(rel: str) -> Path | None:
    """Resolve a brand asset path: absolute as-is, else overlay root then core."""
    path = Path(rel)
    if path.is_absolute():
        return path if path.is_file() else None
    overlay = os.environ.get("KBOTS_OVERLAY", "")
    candidates = ([Path(overlay) / rel] if overlay else []) + [PROJECT_ROOT / rel]
    for candidate in candidates:
        if candidate.is_file():
            return candidate.resolve()
    return None


def _inline_logo(brand: dict, height: int = 40) -> str:
    """Return the brand logo as inline HTML ('' if unavailable).

    The path comes from brand.yaml (operator-controlled config, not agent
    input), so it is read directly. SVG is inlined; raster formats become a
    data URI so the render request-guard never has to fetch a file.
    """
    logo = brand.get("logo") or ""
    if not logo:
        return ""
    path = _resolve_asset(logo)
    if path is None:
        logger.info(f"brand logo not found: {logo}")
        return ""
    if path.suffix.lower() == ".svg":
        svg = path.read_text()
        return f'<div style="height:{height}px">{svg}</div>'
    mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "webp": "image/webp"}.get(path.suffix.lower().lstrip("."), "image/png")
    encoded = base64.b64encode(path.read_bytes()).decode()
    return f'<img src="data:{mime};base64,{encoded}" style="height:{height}px">'


def _is_full_document(html: str) -> bool:
    return html.lstrip().lower().startswith(("<!doctype", "<html"))


_IMG_MIME = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
             ".webp": "image/webp", ".gif": "image/gif", ".svg": "image/svg+xml"}


def _inline_local_images(html: str) -> str:
    """Replace local-file img sources with data URIs.

    Chromium blocks file:// subresources on set_content pages regardless of
    routing, so validated local images (charts, logos) are embedded directly.
    Paths outside the allowed roots are left alone — they render as a broken
    image and the rejection is logged.
    """
    def replace(match: re.Match) -> str:
        src = match.group(2)
        if src.startswith(("data:", "http://", "https://")):
            return match.group(0)
        path_str = src[len("file://"):] if src.startswith("file://") else src
        path = Path(path_str)
        if not path.is_absolute():
            path = path.resolve()  # relative to the service working directory
        mime = _IMG_MIME.get(path.suffix.lower())
        if mime is None:
            return match.group(0)
        if validate_file_path(str(path)) is not None or not path.is_file():
            logger.info(f"design render: not inlining image {src}")
            return match.group(0)
        encoded = base64.b64encode(path.read_bytes()).decode()
        return f'{match.group(1)}data:{mime};base64,{encoded}{match.group(3)}'

    return re.sub(r'(<img[^>]+src=["\'])([^"\']+)(["\'])', replace, html)


def _wrap_html(brand: dict, style: str, body: str, with_chrome: bool = False) -> str:
    """Wrap a fragment in a branded document. with_chrome adds header/footer bands."""
    body = _inline_local_images(body)
    if _is_full_document(body):
        # Full document from the agent — inject brand head after <head> if present
        head_extras = _brand_head(brand, style)
        if "<head>" in body:
            return body.replace("<head>", f"<head>{head_extras}", 1)
        return body.replace("<html>", f"<html><head>{head_extras}</head>", 1) \
            if "<html>" in body else body
    chrome_top = chrome_bottom = ""
    if with_chrome:
        logo = _inline_logo(brand)
        name_block = (f'<div><strong style="font-size:1.05rem">{brand["name"]}</strong>'
                      + (f'<div style="color:var(--brand-muted);font-size:0.8rem">{brand["tagline"]}</div>'
                         if brand.get("tagline") else "") + "</div>")
        chrome_top = (
            '<header style="display:flex;align-items:center;gap:16px;'
            'border-bottom:3px solid var(--brand-primary);padding:8px 0 14px;margin-bottom:28px">'
            f"{logo}{name_block if not logo else ''}"
            "</header>"
        )
        date = datetime.now(timezone.utc).strftime("%Y-%m-%d")
        footer_text = brand.get("footer_text") or brand["name"]
        chrome_bottom = (
            '<footer style="border-top:1px solid var(--brand-muted);color:var(--brand-muted);'
            f'font-size:0.75rem;margin-top:36px;padding-top:10px">{footer_text} · {date}</footer>'
        )
    container = "main" if style == "clean" else "div"
    return (
        "<!doctype html><html><head>"
        + _brand_head(brand, style)
        + "</head><body>"
        + f'<{container} style="padding:24px">{chrome_top}{body}{chrome_bottom}</{container}>'
        + "</body></html>"
    )


# ---------------------------------------------------------------------------
# Chromium rendering
# ---------------------------------------------------------------------------

_dns_cache: dict[str, bool] = {}


def _host_is_public(host: str) -> bool:
    """True if the host resolves only to public addresses. Cached per process."""
    if host in _dns_cache:
        return _dns_cache[host]
    ok = True
    try:
        for info in socket.getaddrinfo(host, None):
            addr = ipaddress.ip_address(info[4][0])
            if any(addr in net for net in _BLOCKED_NETS):
                ok = False
                break
    except (socket.gaierror, ValueError):
        ok = False
    _dns_cache[host] = ok
    return ok


async def _guard_route(route) -> None:
    """Block requests from rendered content to private networks / disallowed files."""
    url = route.request.url
    scheme = url.split(":", 1)[0].lower()
    if scheme in ("data", "about", "blob"):
        await route.continue_()
        return
    if scheme == "file":
        path = url[len("file://"):].split("?")[0]
        if validate_file_path(path) is None:
            await route.continue_()
        else:
            logger.info(f"design render: blocked file url {url}")
            await route.abort()
        return
    if scheme in ("http", "https"):
        host = url.split("/")[2].split(":")[0].split("@")[-1]
        if _host_is_public(host):
            await route.continue_()
        else:
            logger.info(f"design render: blocked private-host url {url}")
            await route.abort()
        return
    await route.abort()


@asynccontextmanager
async def _render_page(html: str, width: int = 1280, height: int = 800, scale: int = 2):
    """Headless Chromium page with the content loaded and requests guarded."""
    from playwright.async_api import async_playwright

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        try:
            page = await browser.new_page(
                viewport={"width": width, "height": height},
                device_scale_factor=scale,
            )
            await page.route("**/*", _guard_route)
            await page.set_content(html, wait_until="networkidle")
            yield page
        finally:
            await browser.close()


def _output_path(ctx: ToolContext, output_file: str, default_name: str) -> Path | str:
    """Resolve an output path (default <workspace>/design/<name>). Path or error."""
    if output_file:
        err = validate_file_path(output_file)
        if err:
            return err
        path = Path(output_file).resolve()
    else:
        base = Path(ctx.project_dir) if ctx.project_dir else KBOTS_TMP
        path = (base / "design" / default_name).resolve()
        err = validate_file_path(str(path))
        if err:
            return err
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        return f"Cannot create output directory {path.parent}: {e}"
    return path


def _timestamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S-%f")


# ---------------------------------------------------------------------------
# Tools
# ---------------------------------------------------------------------------

@tool(name="render_diagram", description="Render a Mermaid diagram to PNG or SVG", category="design")
async def render_diagram(ctx: ToolContext, diagram: str, output_file: str = "",
                         format: str = "png") -> str:
    """Render a diagram from Mermaid text — flowcharts (graph TD), sequence
    diagrams, class/ER diagrams, gantt charts, pie charts, mind maps.

    The diagram is styled with the deployment's brand colors automatically.
    Requires outbound network (mermaid.js from CDN).
    """
    if format not in ("png", "svg"):
        return f"Invalid format '{format}'. Valid: png, svg"
    if not diagram.strip():
        return "Diagram source is empty."

    brand = _load_brand()
    c = brand["colors"]
    theme_vars = {
        "primaryColor": c["primary"], "primaryTextColor": "#FFFFFF",
        "secondaryColor": c["accent"], "tertiaryColor": c["surface"],
        "lineColor": c["secondary"], "textColor": c["text"],
        "fontFamily": brand["font_family"],
    }
    html = f"""<!doctype html><html><head><meta charset="utf-8"></head>
<body style="background:#FFFFFF;margin:0;padding:16px">
<div id="container"></div>
<script type="module">
  import mermaid from "https://cdn.jsdelivr.net/npm/mermaid@11/dist/mermaid.esm.min.mjs";
  mermaid.initialize({{startOnLoad: false, theme: "base",
                       themeVariables: {json.dumps(theme_vars)}}});
  try {{
    const {{svg}} = await mermaid.render("diagram", {json.dumps(diagram)});
    document.getElementById("container").innerHTML = svg;
    window.__result = {{ok: true}};
  }} catch (e) {{
    window.__result = {{ok: false, error: String(e.message || e)}};
  }}
</script></body></html>"""

    out = _output_path(ctx, output_file, f"diagram-{_timestamp()}.{format}")
    if isinstance(out, str):
        return out

    try:
        async with _render_page(html) as page:
            await page.wait_for_function("window.__result !== undefined", timeout=30000)
            result = await page.evaluate("window.__result")
            if not result.get("ok"):
                return f"Mermaid error: {result.get('error', 'unknown')}"
            if format == "svg":
                svg = await page.evaluate("document.getElementById('container').innerHTML")
                out.write_text(svg, encoding="utf-8")
            else:
                await page.locator("#container svg").screenshot(path=str(out))
    except ImportError:
        return _PW_HINT
    except Exception as e:
        return f"Diagram rendering failed: {e}"

    return f"Diagram saved: {out}\nEmbed in reports/HTML with: ![diagram]({out})"


@tool(name="render_svg", description="Save SVG markup and optionally render it to PNG", category="design")
async def render_svg(ctx: ToolContext, svg: str, output_file: str = "",
                     to_png: bool = False) -> str:
    """Save SVG markup you author (graphics, badges, social cards, logos) as a
    .svg file, optionally also rendered to PNG. Use the brand palette for
    consistency — the hex values are listed in this tool's output.
    """
    if not svg.strip():
        return "SVG content is empty."
    try:
        root = ET.fromstring(svg)
    except ET.ParseError as e:
        return f"SVG is not valid XML: {e}"
    if not root.tag.endswith("svg"):
        return "Content does not look like an <svg> document."

    out = _output_path(ctx, output_file, f"graphic-{_timestamp()}.svg")
    if isinstance(out, str):
        return out
    if out.suffix.lower() != ".svg":
        out = out.with_suffix(".svg")
    try:
        out.write_text(svg, encoding="utf-8")
    except OSError as e:
        return f"Failed to write SVG: {e}"

    lines = [f"SVG saved: {out}"]
    if to_png:
        png_path = out.with_suffix(".png")
        html = f'<!doctype html><html><body style="margin:0;background:transparent">{svg}</body></html>'
        try:
            async with _render_page(html) as page:
                await page.locator("svg").first.screenshot(
                    path=str(png_path), omit_background=True)
            lines.append(f"PNG saved: {png_path}")
        except ImportError:
            lines.append(f"PNG skipped — {_PW_HINT}")
        except Exception as e:
            lines.append(f"PNG rendering failed: {e}")
    lines.append(_brand_tokens_line(_load_brand()))
    return "\n".join(lines)


@tool(name="render_html", description="Render HTML to a PNG screenshot (mockups, landing pages, emails)",
      category="design")
async def render_html(ctx: ToolContext, html: str, output_file: str = "",
                      style: str = "clean", width: int = 1280,
                      full_page: bool = True) -> str:
    """Render HTML to a PNG image — for UI mockups, landing pages, email
    templates, social cards.

    style: 'clean' (semantic HTML auto-styled via Pico.css + brand),
    'tailwind' (write Tailwind utility classes; brand tokens available as
    text-brand, bg-brand-accent, …), or 'none' (your own CSS; brand CSS
    variables like var(--brand-primary) still work).
    Pass a fragment (body content) or a full <html> document.
    """
    if style not in HTML_STYLES:
        return f"Invalid style '{style}'. Valid: {', '.join(HTML_STYLES)}"
    if not html.strip():
        return "HTML content is empty."
    width = max(320, min(width, 3840))

    brand = _load_brand()
    doc = _wrap_html(brand, style, html, with_chrome=False)
    out = _output_path(ctx, output_file, f"mockup-{_timestamp()}.png")
    if isinstance(out, str):
        return out

    try:
        async with _render_page(doc, width=width) as page:
            await page.screenshot(path=str(out), full_page=full_page)
    except ImportError:
        return _PW_HINT
    except Exception as e:
        return f"HTML rendering failed: {e}"

    return f"Screenshot saved: {out}\n{_brand_tokens_line(brand)}"


@tool(name="html_to_pdf", description="Render HTML to a designed PDF (reports, one-pagers)", category="design")
async def html_to_pdf(ctx: ToolContext, html: str, output_file: str = "",
                      style: str = "clean", page_size: str = "A4",
                      landscape: bool = False) -> str:
    """Render HTML to a print-quality PDF via Chromium — full modern CSS
    (grid, flexbox, @page rules, web fonts, embedded local images).

    Use this instead of write_report when layout and branding matter: pass a
    fragment and get a branded document (logo header band + footer added
    automatically), or pass a full <html> document for total control.
    Charts from create_chart can be embedded with <img src="/abs/path.png">
    — local images inside the allowed directories are inlined automatically.
    """
    if style not in HTML_STYLES:
        return f"Invalid style '{style}'. Valid: {', '.join(HTML_STYLES)}"
    if page_size not in PAGE_SIZES:
        return f"Invalid page_size '{page_size}'. Valid: {', '.join(PAGE_SIZES)}"
    if not html.strip():
        return "HTML content is empty."

    brand = _load_brand()
    doc = _wrap_html(brand, style, html, with_chrome=True)
    out = _output_path(ctx, output_file, f"document-{_timestamp()}.pdf")
    if isinstance(out, str):
        return out

    try:
        async with _render_page(doc, width=1080) as page:
            await page.pdf(path=str(out), format=page_size, landscape=landscape,
                           print_background=True,
                           margin={"top": "14mm", "bottom": "14mm",
                                   "left": "12mm", "right": "12mm"})
    except ImportError:
        return _PW_HINT
    except Exception as e:
        return f"PDF rendering failed: {e}"

    size_kb = out.stat().st_size / 1024
    return f"PDF saved: {out} ({size_kb:.1f} KB)\n{_brand_tokens_line(brand)}"


@tool(name="create_slides", description="Create a PowerPoint deck from a markdown outline", category="design")
async def create_slides(ctx: ToolContext, title: str, content: str,
                        output_file: str = "") -> str:
    """Create a .pptx slide deck from a markdown outline.

    Format: slides separated by lines containing only '---'.
    Per slide: '# Heading' is the slide title, '- item' bullets (indent two
    spaces per nesting level), plain lines become body text, '![alt](path)'
    adds an image. A slide with only a title (and optional one paragraph)
    becomes a title slide.
    If <overlay>/config/brand.pptx exists it is used as the template, so the
    deck inherits your corporate slide master.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
    except ImportError:
        return _PPTX_HINT

    if not content.strip():
        return "Slide content is empty."

    brand = _load_brand()
    primary = RGBColor.from_string(brand["colors"]["primary"].lstrip("#"))

    template = resolve_config_file("brand.pptx")
    prs = Presentation(str(template)) if template.is_file() else Presentation()

    raw_slides = [s.strip() for s in re.split(r"^---\s*$", content, flags=re.M) if s.strip()]
    if not raw_slides:
        return "No slides found in the outline."

    img_pattern = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
    made = 0
    skipped_images: list[str] = []

    for index, block in enumerate(raw_slides):
        slide_title = ""
        bullets: list[tuple[int, str]] = []
        paragraphs: list[str] = []
        images: list[str] = []

        for line in block.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            m = img_pattern.match(stripped)
            if m:
                img = m.group(1)
                if validate_file_path(img) is None and Path(img).is_file():
                    images.append(img)
                else:
                    skipped_images.append(img)
                continue
            if stripped.startswith("#"):
                slide_title = stripped.lstrip("#").strip()
            elif stripped.startswith(("- ", "* ")):
                indent = (len(line) - len(line.lstrip())) // 2
                bullets.append((min(indent, 4), stripped[2:].strip()))
            else:
                paragraphs.append(stripped)

        is_title_slide = index == 0 and not bullets and not images
        layout = prs.slide_layouts[0 if is_title_slide else 1]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = slide_title or (title if is_title_slide else f"Slide {index + 1}")
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = primary

        # Body placeholder: subtitle on title layout, content on others
        body = next((ph for ph in slide.placeholders if ph.placeholder_format.idx != 0), None)
        body_lines: list[tuple[int, str]] = [(0, p) for p in paragraphs] + bullets
        if body is not None and body_lines:
            tf = body.text_frame
            first = True
            for level, text_line in body_lines:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                para.text = text_line
                para.level = level
                first = False

        for i, img in enumerate(images):
            try:
                slide.shapes.add_picture(img, Inches(5.2), Inches(1.6 + i * 2.6), width=Inches(4.3))
            except Exception as e:
                skipped_images.append(f"{img} ({e})")
        made += 1

    slug = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-") or "slides"
    out = _output_path(ctx, output_file, f"{slug}-{_timestamp()}.pptx")
    if isinstance(out, str):
        return out
    if out.suffix.lower() != ".pptx":
        out = out.with_suffix(".pptx")
    try:
        prs.save(str(out))
    except OSError as e:
        return f"Failed to save deck: {e}"

    msg = f"Deck saved: {out} ({made} slides)"
    if skipped_images:
        msg += "\nSkipped images: " + ", ".join(skipped_images)
    return msg
